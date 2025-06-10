#!/usr/bin/env python3
import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import base64
from io import BytesIO

def extract_ppt_content(ppt_path):
    """Extract text and images from PowerPoint presentation"""
    prs = Presentation(ppt_path)
    slides_data = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            'slide_number': slide_num,
            'title': '',
            'content': [],
            'images': []
        }
        
        # Extract text content
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Check if it's likely a title (first text shape or largest)
                if not slide_data['title'] and shape.text.strip():
                    slide_data['title'] = shape.text.strip()
                else:
                    slide_data['content'].append(shape.text.strip())
            
            # Extract images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    # Save image with slide and shape info
                    image_filename = f"slide_{slide_num}_image_{len(slide_data['images']) + 1}.{image.ext}"
                    image_path = f"extracted_images/{image_filename}"
                    
                    # Create directory if it doesn't exist
                    os.makedirs('extracted_images', exist_ok=True)
                    
                    # Save image file
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                    
                    slide_data['images'].append({
                        'filename': image_filename,
                        'path': image_path,
                        'format': image.ext
                    })
                except Exception as e:
                    print(f"Error extracting image from slide {slide_num}: {e}")
        
        slides_data.append(slide_data)
    
    return slides_data

def main():
    ppt_path = "attached_assets/PPT new_1749576426036.pptx"
    
    if not os.path.exists(ppt_path):
        print(f"PowerPoint file not found: {ppt_path}")
        return
    
    try:
        slides_content = extract_ppt_content(ppt_path)
        
        # Save extracted content to JSON
        with open('ppt_content.json', 'w', encoding='utf-8') as f:
            json.dump(slides_content, f, indent=2, ensure_ascii=False)
        
        print(f"Successfully extracted content from {len(slides_content)} slides")
        print("Content saved to ppt_content.json")
        
        # Print summary
        for slide in slides_content:
            print(f"\nSlide {slide['slide_number']}: {slide['title'][:50]}...")
            if slide['images']:
                print(f"  - {len(slide['images'])} images extracted")
            print(f"  - {len(slide['content'])} text blocks")
            
    except Exception as e:
        print(f"Error processing PowerPoint: {e}")

if __name__ == "__main__":
    main()